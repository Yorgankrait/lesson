# -*- coding: utf-8 -*-
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required, user_passes_test
from django.http import HttpResponse, JsonResponse
from django.contrib.auth import login, authenticate
from .models import Lesson, Attendance, ChatMessage, Student, UnknownQuestion, News, AboutPage, TeacherResource, UserIdea, TeacherArticle
from .forms import LessonForm, CustomUserCreationForm, ExcelUploadForm, AttendanceForm, StudentForm, TeacherArticleForm, TeacherResourceForm
import pandas as pd
import os
from datetime import datetime, timedelta
from django.utils.dateformat import format
from itertools import groupby
from operator import attrgetter
from django.views.decorators.csrf import csrf_exempt
import json
from django.conf import settings
from pptx import Presentation
import io
import base64
from PIL import Image, ImageDraw, ImageFont
import sys
from io import StringIO
from django.core.exceptions import ValidationError
from django.utils import timezone
from django.contrib import messages
from django.views.decorators.http import require_http_methods
from django.contrib.sessions.models import Session
from django.contrib.sessions.backends.db import SessionStore
import logging
from django.contrib.auth.hashers import make_password
from django.contrib.auth.models import User
from .forms import PasswordResetRequestForm
from django.core.exceptions import PermissionDenied
from django.urls import reverse
from django.core.files.storage import default_storage

logger = logging.getLogger(__name__)


KESHA_GREETING = """Я могу рассказать о любом разделе сайта и правах доступа. Доступные команды:
- "главная" - информация о главной странице
- "ученики" - информация о разделе учеников и проектов
- "материалы" - информация об учебных материалах 
- "интерпретатор" - информация о Python интерпретаторе
- "посещаемость" - информация о посещаемости
- "права" - информация о правах разных типов пользователей
- "предложить идею" - поделиться идеей по улучшению сайта
- "keshagpt" - информация о продвинутом ассистенте (для администраторов и активированных учителей)"""

KESHA_RIGHTS = """На сайте есть несколько типов пользователей с разными правами:

Все пользователи:
- Доступ к Python интерпретатору
- Просмотр информации об авторе
- Базовое общение с чат-ботом

Авторизованные пользователи:
- Расширенное общение с чат-ботом
- Возможность предлагать идеи
- Просмотр учеников и проектов
- Просмотр посещаемости

Ученики (после активации):
- Доступ к учебным материалам
- Возможность скачивания материалов

Родители:
- Нет доступа к учебным материалам

Учителя (после активации):
- Доступ к учебным материалам
- Доступ к разделу "Полезное для учителей"
- Доступ к KeshaGPT

Администраторы:
- Управление пользователями
- Активация учетных записей
- Управление материалами
- Обработка вопросов и идей
- Полный доступ к KeshaGPT"""

KESHA_KESHAGPT = """KeshaGPT - это продвинутый ассистент на базе GPT, дступный только для администраторов и активированных учителей.

Возможности KeshaGPT:
- Помощь в подготовке учебных материалов
- Ответы на сложные вопросы по программированию
- Генерация примеров кода и объяснений
- Проверка кода учеников
- Рекомендации по улучшению материалов

Для получения доступа к KeshaGPT учителям необходимо:
1. Зарегистрироваться как учитель
2. Дождаться активации аккаунта администратором
3. Получить специальный доступ к KeshaGPT"""

def fix_session(request):
    if not request.session.exists(request.session.session_key):
        request.session.create()
    try:
        Session.objects.get(session_key=request.session.session_key)
    except Session.DoesNotExist:
        request.session.save()
    return request

def wrap_view(view_func):
    def wrapped(request, *args, **kwargs):
        request = fix_session(request)
        return view_func(request, *args, **kwargs)
    return wrapped

@wrap_view
def home(request):
    news = News.objects.filter(is_published=True).order_by('-created_at')
    return render(request, 'home.html', {'news': news})

def check_student_access(user):
    if not user.is_authenticated:
        return False
    profile = user.userprofile
    return profile.user_type == 'student' and profile.is_student_activated

def check_parent_access(user):
    if not user.is_authenticated:
        return False
    profile = user.userprofile
    return profile.user_type == 'parent'

def check_teacher_access(user):
    if not user.is_authenticated:
        return False
    profile = user.userprofile
    return profile.user_type == 'teacher' and profile.is_teacher_activated

def check_teacher_or_admin_access(user):
    if not user.is_authenticated:
        return False
    if user.is_superuser:  # Проверяем, является ли пользователь админом
        return True
    profile = user.userprofile
    return profile.user_type == 'teacher' and profile.is_teacher_activated

@wrap_view
def students(request):
    if not request.user.is_authenticated:
        messages.warning(request, 'Для просмотра этой страницы необходимо войти в систему.')
        return redirect('login')
    
    active_students = Student.objects.filter(active=True).select_related('project')
    students_data = []
    
    for student in active_students:
        student_data = {
            'name': student.name,
            'project': student.project.title if student.project else 'Проект не назначен',
            'description': student.project.description if student.project else 'Описание отсутствует'
        }
        students_data.append(student_data)
    
    return render(request, 'students.html', {'students': students_data})

@wrap_view
def lessons(request):
    if not request.user.is_authenticated:
        messages.warning(request, 'Для просмотра учебных материалов необходимо войти в систему.')
        return redirect('login')
    
    # Проверяем, является ли пользователь родителем
    if request.user.userprofile.user_type == 'parent':
        messages.warning(request, 'У родителей нет доступа к учебным материалам.')
        return render(request, 'lessons.html', {'error_message': 'У родитлей нет доступа к учебным мариалам.'})
    
    # Проверяем, является ли пользователь неактивированным учителем
    if request.user.userprofile.user_type == 'teacher' and not request.user.userprofile.is_teacher_activated:
        messages.warning(request, 'Ваша учетная запись учителя ожидает активации администратором. После активации вам станут доступны учебные материалы.')
        return render(request, 'lessons.html', {'error_message': 'Ваша учетная запись учителя ожидает активации администратором. После активации вам станут доступны учебные материалы.'})
    
    # Проверяем, является ли пользователь неактивированным учеником
    if request.user.userprofile.user_type == 'student' and not request.user.userprofile.is_student_activated:
        messages.warning(request, 'Ваша учетная запись ученика ожидает активации администратором. После активации вам станут доступны учебные материалы. Пока вы можете пользоваться Python интерпретатором, просматривать информацию о посещаемости, а также посмотреть информацию об авторе проекта.')
        return render(request, 'lessons.html', {'error_message': 'Ваша учетная запись ученика ожидает активации администратором. После активации вам станут доступны учебные материалы. Пока вы можете пользоваться Python интерпретатором, просматривать информацию о посещаемости, а также посмотреть информацию об авторе проекта.'})
    
    lessons = Lesson.objects.all().order_by('-upload_date')
    return render(request, 'lessons.html', {'lessons': lessons})

@csrf_exempt
@wrap_view
def chat_message(request):
    if not request.user.is_authenticated:
        return JsonResponse({
            'status': 'success', 
            'message': ('Здравствуйте! Я помогу вам разобраться с сайтом.\n\n'
                       'Сейчас вам доступны:\n'
                       '- Python Интерпретатор\n'
                       '- Просмотр информации "Обо мне"\n\n'
                       'Чтобы получить полный доступ ко всем функциям сайта, пожалуйста, '
                       'зарегистрируйтесь или войдите в систему.\n\n'
                       'Для регистрации нажмите кнопку "Регистрация" внизу страницы.\n'
                       'Для входа нажмите "Вход для пользователей".')
        })
    
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            message = data.get('message').lower()
            
            # Сохраняем сообщее пользователя
            ChatMessage.objects.create(user=request.user, message=message, is_bot=False)
            
            # Специальная команда для администратора
            if request.user.is_superuser and message == 'что нового':
                # Получаем непросмотренные идеи
                unreviewed_ideas = UserIdea.objects.filter(is_reviewed=False).order_by('-created_at')
                
                if unreviewed_ideas.exists():
                    bot_response = 'Новые идеи от пользователей:\n\n'
                    for idx, idea in enumerate(unreviewed_ideas, 1):
                        bot_response += f"{idx}. От пользователя {idea.user.username} ({idea.created_at.strftime('%d.%m.%Y %H:%M')}):\n{idea.idea}\n\n"
                else:
                    bot_response = 'Новых непросмотренных идей пока нет.'
                
                ChatMessage.objects.create(user=request.user, message=bot_response, is_bot=True)
                return JsonResponse({'status': 'success', 'message': bot_response})
            
            # Проверяем, ожидает л�� пользователь ввода идеи
            if 'waiting_for_idea' in request.session:
                # Создаем новую идею
                UserIdea.objects.create(
                    user=request.user,
                    idea=message
                )
                # Удаляем флаг ожидания
                del request.session['waiting_for_idea']
                bot_response = 'Спасибо за вашу идею! Я передал её администратору.'
                ChatMessage.objects.create(user=request.user, message=bot_response, is_bot=True)
                return JsonResponse({'status': 'success', 'message': bot_response})
            
            # Определяем словарь responses
            responses = {
                # Приветствия и прощания
                'привет': 'Здравствуйте! Я помогу вам разобраться с сайтом. Спросите меня о любом разделе, напишите "права" чтобы узнать о возможностях разных типов пользователей, или "предложить идею" чтобы поделиться своими идеями по улучшению сайта.',
                'пока': 'До свидания! Буду рад помочь снова!',
                'как дела': 'У меня всё хорошо, готов помогать! Что вас интересует? Напишите "права" чтобы узнать о возможностях разных типов пользователей, или "предложить идею" тобы поделиться своими идеями.',
                
                # Описание разделов
                'главная': 'На главной странице вы найдете общую информацию о нашем обучающем проекте и меня - вашего помощника по сайту.',
                'ученики': 'В аделе "Ученики и их проекты" вы можете увидеть работы наших учеников, их достижения и проекты, которые они создали во время обучения.',
                'материалы': 'В разделе "Учебные материалы" собраны все презентации и уроки по программированию.',
                'интерпретатор': 'Python Интерпретатор позволяет писать и выполнять код прямо в браузере.',
                
                # Дополнительные функции
                'посещаемость': 'В разделе "Посещаемость" можно отслеживать присутствие учеников на занятиях.',
                'регистрация': 'При регистрации вы можете вбрать тип учетной записи: ученик, родитель или учитель. У каждого типа свои права доступа (напишите "права" для подробностей).',
                'вход': 'Для входа а сайт используйте кнопку "Вход для пользователей" в нижней части страницы.',
                'админ': 'Вход для администратора доступен по соответствующей кнопке в утере сайта.',
                
                # Помощь
                'помощь': KESHA_GREETING,
                'права': KESHA_RIGHTS,
                'keshagpt': KESHA_KESHAGPT,
                
                # Команда для предложения идеи
                'предложить идею': 'Конечно, я передам вашу идею администратору, опишите её',
                
                # обавляем информацию о команде для администратора
                'команды админа': ('Доступные команды для администратора:\n'
                                 '- "что нового" - просмотр непросмотренных идей пользователей'),
            }
            
            # Сначала проверяем, есть ли ответ в базе отвеченных вопросов
            answered_question = UnknownQuestion.objects.filter(
                question__iexact=message,
                is_answered=True,
                answer__isnull=False
            ).first()
            
            if answered_question:
                bot_response = answered_question.answer
            else:
                # Поиск отета в стандартнх ответах
                bot_response = None
                for key in responses:
                    if key in message:
                        bot_response = responses[key]
                        if key == 'предложить идею':
                            request.session['waiting_for_idea'] = True
                        break
                
                # Если ответ не найден, сохраняем вопрос в базе
                if bot_response is None:
                    UnknownQuestion.objects.create(
                        user=request.user,
                        question=message,
                        is_answered=False
                    )
                    bot_response = ('Извините, я пока не знаю ответа на этот вопрос. '
                                  'Я передал ваш вопрос администратору. '
                                  'Попробуйте спросить о конкретном разделе сайта '
                                  'или напшите "помощь" для получения списка доступных коанд.')
            
            # Сохраняем ответ бота
            ChatMessage.objects.create(user=request.user, message=bot_response, is_bot=True)
            
            return JsonResponse({'status': 'success', 'message': bot_response})
            
        except Exception as e:
            logger.error(f"Error in chat_message: {str(e)}")
            return JsonResponse({'status': 'error', 'message': str(e)}, status=500)
    
    elif request.method == 'GET':
        if not request.user.is_authenticated:
            return JsonResponse({'status': 'success', 'messages': []})
        messages = ChatMessage.objects.filter(user=request.user).values('message', 'is_bot')
        return JsonResponse({'status': 'success', 'messages': list(messages)})
    
    elif request.method == 'DELETE':
        if not request.user.is_authenticated:
            return JsonResponse({'status': 'error', 'message': 'Authentication required'}, status=401)
        ChatMessage.objects.filter(user=request.user).delete()
        return JsonResponse({'status': 'success', 'message': 'Chat history cleared'})
    
    return JsonResponse({'status': 'error', 'message': 'Invalid request method'}, status=400)

@login_required
@wrap_view
def view_lesson(request, lesson_id):
    lesson = get_object_or_404(Lesson, id=lesson_id)
    
    try:
        pptx_path = lesson.presentation.path
        prs = Presentation(pptx_path)
        
        slides = []
        for slide in prs.slides:
            img = Image.new('RGB', (960, 540), (0, 0, 0))  # Черный фон
            draw = ImageDraw.Draw(img)
            
            try:
                # Определяем систему и выбираем соответствующе шрифты
                if os.name == 'nt':  # Windows
                    title_font = ImageFont.truetype("C:\\Windows\\Fonts\\Arial.ttf", 32)
                    content_font = ImageFont.truetype("C:\\Windows\\Fonts\\Arial.ttf", 24)
                    code_font = ImageFont.truetype("C:\\Windows\\Fonts\\Consola.ttf", 20)
                else:  # macOS/Linux
                    title_font = ImageFont.truetype("/Library/Fonts/Arial Bold.ttf", 32)
                    content_font = ImageFont.truetype("/Library/Fonts/Arial.ttf", 24)
                    code_font = ImageFont.truetype("/Library/Fonts/Courier New.ttf", 20)
            except IOError as e:
                logger.error(f"Font loading error: {str(e)}")
                # Если не удалось загрузить системные шрфты, используем дефолтный
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
                code_font = ImageFont.load_default()
            
            y_position = 20
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text = shape.text.strip()
                    if text:
                        # Явно указываем кодировку для Windows
                        if os.name == 'nt':
                            text = text.encode('utf-8').decode('utf-8')
                            
                        if shape.name == 'Title':
                            text_width = draw.textbbox((0, 0), text, font=title_font)[2]
                            x_position = (960 - text_width) / 2
                            draw.text((x_position, y_position), text, fill=(255, 255, 255), font=title_font)
                            y_position += 60
                        else:
                            paragraphs = text.split('\n\n')
                            for paragraph in paragraphs:
                                if paragraph.strip().startswith('# ') or '=' in paragraph or 'print(' in paragraph:
                                    lines = paragraph.split('\n')
                                    for line in lines:
                                        words = line.split()
                                        line = []
                                        for word in words:
                                            line.append(word)
                                            bbox = draw.textbbox((0, 0), ' '.join(line), font=code_font)
                                            w = bbox[2] - bbox[0]
                                            if w > 880:
                                                draw.text((40, y_position), ' '.join(line[:-1]), fill=(0, 255, 0), font=code_font)
                                                y_position += 25
                                                line = [word]
                                        if line:
                                            draw.text((40, y_position), ' '.join(line), fill=(0, 255, 0), font=code_font)
                                            y_position += 25
                                else:
                                    words = paragraph.split()
                                    line = []
                                    for word in words:
                                        line.append(word)
                                        bbox = draw.textbbox((0, 0), ' '.join(line), font=content_font)
                                        w = bbox[2] - bbox[0]
                                        if w > 880:
                                            draw.text((40, y_position), ' '.join(line[:-1]), fill=(255, 255, 255), font=content_font)
                                            y_position += 30
                                            line = [word]
                                    if line:
                                        draw.text((40, y_position), ' '.join(line), fill=(255, 255, 255), font=content_font)
                                        y_position += 30
                                y_position += 20
                
                if y_position > 520:
                    break
            
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            img_str = base64.b64encode(buffered.getvalue()).decode()
            slides.append(f"data:image/png;base64,{img_str}")
        
        context = {
            'lesson': lesson,
            'slides': slides,
        }
        
        return render(request, 'view_lesson.html', context)
    
    except Exception as e:
        logger.error(f'Error in view_lesson: {str(e)}')
        error_message = f'Ошибка при открытии презентации: {str(e)}'
        return render(request, 'error.html', {'error_message': error_message})

@login_required
@wrap_view
def download_lesson(request, lesson_id):
    lesson = get_object_or_404(Lesson, id=lesson_id)
    response = HttpResponse(lesson.presentation, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = f'attachment; filename="{lesson.title}.pptx"'
    return response

@login_required
@wrap_view
def upload_lesson(request):
    if not request.user.is_staff:
        return redirect('lessons')
    if request.method == 'POST':
        form = LessonForm(request.POST, request.FILES)
        if form.is_valid():
            presentation = request.FILES['presentation']
            if not presentation.name.lower().endswith('.pptx'):
                form.add_error('presentation', 'Загруженный файл не является PowerPoint презентацией (.pptx).')
            else:
                form.save()
                return redirect('lessons')
    else:
        form = LessonForm()
    return render(request, 'upload_lesson.html', {'form': form})

@user_passes_test(lambda u: u.is_superuser)
def delete_lesson(request, lesson_id):
    lesson = get_object_or_404(Lesson, pk=lesson_id)
    if request.method == 'POST':
        # Удаляем файл презентации
        if lesson.presentation:
            if os.path.exists(lesson.presentation.path):
                os.remove(lesson.presentation.path)
        # Удаляем урок
        lesson.delete()
        messages.success(request, 'Урок успешно удален')
        return redirect('lessons')
    return render(request, 'delete_lesson.html', {'lesson': lesson})

@wrap_view
def register(request):
    if request.method == 'POST':
        form = CustomUserCreationForm(request.POST)
        if form.is_valid():
            user = form.save()
            login(request, user)
            return redirect('home')
    else:
        form = CustomUserCreationForm()
    return render(request, 'register.html', {'form': form})

@wrap_view
def python_interpreter(request):
    # Без проверки авторизации
    return render(request, 'python_interpreter.html')

@csrf_exempt
@wrap_view
def run_python_code(request):
    # Без проверки авторизации
    if request.method == 'POST':
        code = json.loads(request.body)['code']
        
        old_stdout = sys.stdout
        old_stderr = sys.stderr
        redirected_output = sys.stdout = StringIO()
        redirected_error = sys.stderr = StringIO()

        try:
            exec(code)
            output = redirected_output.getvalue()
            error = redirected_error.getvalue()
            if error:
                output = f"Error:\n{error}\n\nOutput:\n{output}"
        except Exception as e:
            output = str(e)
        finally:
            sys.stdout = old_stdout
            sys.stderr = old_stderr

        return JsonResponse({'output': output})

@user_passes_test(lambda u: u.is_staff)
@wrap_view
def add_student(request):
    if request.method == 'POST':
        form = StudentForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, 'Ученик успешно добавлен.')
            return redirect('attendance')
    else:
        form = StudentForm()
    return render(request, 'add_student.html', {'form': form})

@wrap_view
def attendance(request):
    if not request.user.is_authenticated:
        messages.warning(request, 'Для просмотра псещаемости необходимо войти в систему.')
        return redirect('login')
    
    attendances = Attendance.objects.all().order_by('-date', 'student')
    
    grouped_attendances = []
    for date, group in groupby(attendances, key=lambda x: x.date.strftime('%B %Y')):
        grouped_attendances.append((date, list(group)))
    
    context = {
        'grouped_attendances': grouped_attendances,
    }
    if request.user.is_staff:
        context['students'] = Student.objects.all()
    return render(request, 'attendance.html', context)

@user_passes_test(lambda u: u.is_staff)
@wrap_view
def mark_attendance(request):
    if request.method == 'POST':
        form = AttendanceForm(request.POST)
        if form.is_valid():
            date = form.cleaned_data['date']
            no_class = form.cleaned_data['no_class']
            present_students = form.cleaned_data['students']
            
            if no_class:
                for student in Student.objects.all():
                    Attendance.objects.update_or_create(
                        student=student,
                        date=date,
                        defaults={'present': False, 'no_class': True}
                    )
            else:
                for student in Student.objects.all():
                    Attendance.objects.update_or_create(
                        student=student,
                        date=date,
                        defaults={'present': student in present_students, 'no_class': False}
                    )
            
            messages.success(request, 'Посещаемость успешно отмечена.')
            return redirect('attendance')
    else:
        form = AttendanceForm()
    return render(request, 'mark_attendance.html', {'form': form})

@require_http_methods(["GET", "POST"])
@wrap_view
def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')
        user = authenticate(request, username=username, password=password)
        if user is not None:
            login(request, user)
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({'success': True, 'redirect_url': '/'})
            return redirect('home')
        elif request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            return JsonResponse({'success': False, 'message': 'Invalid credentials'})
        else:
            messages.error(request, 'Invalid username or password')
    return render(request, 'login.html')

def password_reset_request(request):
    if request.method == 'POST':
        form = PasswordResetRequestForm(request.POST)
        if form.is_valid():
            username = form.cleaned_data['username']
            email = form.cleaned_data['email']
            try:
                user = User.objects.get(username=username, email=email)
                
                # Если пользоватль найден и еть новый пароль в запросе
                if 'new_password1' in request.POST and 'new_password2' in request.POST:
                    new_password1 = request.POST['new_password1']
                    new_password2 = request.POST['new_password2']
                    
                    if new_password1 != new_password2:
                        messages.error(request, 'Пароли не совпадают!')
                        return render(request, 'password_reset_form.html', {'user_found': True})
                    
                    if len(new_password1) < 8:
                        messages.error(request, 'Пароль должен содержать минимум 8 символо!')
                        return render(request, 'password_reset_form.html', {'user_found': True})
                    
                    try:
                        # Устанавливаем новый пароль
                        user.set_password(new_password1)
                        user.save()
                        
                        # Авторзуем пользователя
                        auth_user = authenticate(username=username, password=new_password1)
                        if auth_user is not None:
                            login(request, auth_user)
                            messages.success(request, 'Пароль успешно изменен! Сейчас вы будете перенаправлены на главную страницу.')
                            return render(request, 'password_reset_success.html')
                        else:
                            messages.warning(request, 'Пароль изменен, но автоматическая воризация не удалас. Пожалуйста, войдите с новым паролем.')
                            return redirect('login')
                    except Exception as e:
                        messages.error(request, f'Ошибка при смене пароя: {str(e)}')
                        return render(request, 'password_reset_form.html', {'user_found': True})
                
                # Если пльзователь найден, но это первый шаг (проверка логина и email)
                return render(request, 'password_reset_form.html', {
                    'user_found': True,
                    'username': username,
                    'email': email
                })
                
            except User.DoesNotExist:
                messages.error(request, 'Пользователь с такими анными не найден')
                return render(request, 'password_reset_request.html', {'form': form})
    else:
        form = PasswordResetRequestForm()
    
    return render(request, 'password_reset_request.html', {'form': form})

@wrap_view
def about(request):
    # Доступно всем авторизованным пользователям
    about_info = AboutPage.objects.first()  # Берем первую запись
    return render(request, 'about.html', {'about_info': about_info})

@wrap_view
def teacher_resources(request):
    if not request.user.is_authenticated:
        return redirect('login')
    
    if not (request.user.is_superuser or 
            (hasattr(request.user, 'userprofile') and 
             request.user.userprofile.user_type == 'teacher' and 
             request.user.userprofile.is_teacher_activated)):
        raise PermissionDenied
    
    articles = TeacherArticle.objects.filter(is_published=True).prefetch_related('resources')
    return render(request, 'teacher_resources.html', {'articles': articles})

def keshagpt_view(request, feature=None):
    # Проверяем, является ли пользователь админом или активированным учителем
    if not request.user.is_authenticated:
        raise PermissionDenied
    
    is_authorized = (
        request.user.is_superuser or 
        (hasattr(request.user, 'userprofile') and 
         request.user.userprofile.user_type == 'teacher' and 
         request.user.userprofile.is_teacher_activated)
    )
    
    if not is_authorized:
        raise PermissionDenied
    
    # Базовый URL blackbox.ai
    base_url = 'https://www.blackbox.ai'
    
    # Словарь соответствия функций и их URL
    feature_urls = {
        'chat': base_url,
        'image': 'https://www.blackbox.ai/agent/ImageGenerationLV45LJp'
    }
    
    # Определяем URL для рендеринга
    url = feature_urls.get(feature, base_url)
    
    return render(request, 'keshagpt.html', {'url': url})

def is_teacher(user):
    return user.groups.filter(name='Teacher').exists() or user.is_superuser

@user_passes_test(lambda u: u.is_superuser or (hasattr(u, 'userprofile') and u.userprofile.user_type == 'teacher' and u.userprofile.is_teacher_activated))
def teacher_articles_list(request):
    articles = TeacherArticle.objects.filter(is_published=True).order_by('-created_at')
    context = {
        'articles': articles,
        'is_superuser': request.user.is_superuser  # Явно передаем флаг суперпользователя
    }
    return render(request, 'main_app/teacher_articles_list.html', context)

@user_passes_test(lambda u: u.is_superuser)  # Только админ может создавать статьи
def create_teacher_article(request):
    if request.method == 'POST':
        form = TeacherArticleForm(request.POST, request.FILES)
        if form.is_valid():
            article = form.save()
            
            # Собираем контент из блоков
            content = []
            content_types = request.POST.getlist('content_type[]', [])
            contents = request.POST.getlist('content[]', [])
            images = request.FILES.getlist('images[]', [])
            
            image_index = 0
            text_index = 0
            
            for content_type in content_types:
                if content_type == 'text':
                    if text_index < len(contents) and contents[text_index].strip():
                        content.append({
                            'type': 'text',
                            'content': contents[text_index]
                        })
                    text_index += 1
                elif content_type == 'image':
                    if image_index < len(images):
                        image_name = f"article_{article.id}_image_{image_index}{os.path.splitext(images[image_index].name)[1]}"
                        image_path = default_storage.save(f'article_images/{image_name}', images[image_index])
                        content.append({
                            'type': 'image',
                            'path': image_path
                        })
                        image_index += 1
            
            article.content = json.dumps(content)
            article.save()
            
            messages.success(request, 'Статья успешно создана')
            return redirect('teacher_article_detail', pk=article.pk)
    else:
        form = TeacherArticleForm()
    
    return render(request, 'main_app/create_teacher_article.html', {'form': form})

@user_passes_test(lambda u: u.is_superuser or (hasattr(u, 'userprofile') and u.userprofile.user_type == 'teacher' and u.userprofile.is_teacher_activated))
def teacher_article_detail(request, pk):
    article = get_object_or_404(TeacherArticle, pk=pk)
    resources = article.resources.all().order_by('order')
    resource_form = TeacherResourceForm() if request.user.is_superuser else None
    
    if request.method == 'POST' and request.user.is_superuser:
        resource_form = TeacherResourceForm(request.POST, request.FILES)
        if resource_form.is_valid():
            resource = resource_form.save(commit=False)
            resource.article = article
            resource.save()
            messages.success(request, 'Файл успешно добавлен')
            return redirect('teacher_article_detail', pk=pk)
            
    return render(request, 'main_app/teacher_article_detail.html', {
        'article': article,
        'resources': resources,
        'resource_form': resource_form
    })

@user_passes_test(lambda u: u.is_superuser)
def edit_teacher_article(request, pk):
    article = get_object_or_404(TeacherArticle, pk=pk)
    if request.method == 'POST':
        form = TeacherArticleForm(request.POST, instance=article)
        if form.is_valid():
            article = form.save()
            messages.success(request, 'Статья успешно обновлена')
            return redirect('teacher_article_detail', pk=article.pk)
    else:
        form = TeacherArticleForm(instance=article)
    return render(request, 'main_app/edit_teacher_article.html', {'form': form, 'article': article})

@user_passes_test(lambda u: u.is_superuser)
def delete_teacher_article(request, pk):
    article = get_object_or_404(TeacherArticle, pk=pk)
    if request.method == 'POST':
        # Удаляем превью изображение, если оно есть
        if article.preview_image:
            if os.path.exists(article.preview_image.path):
                os.remove(article.preview_image.path)

        # Получаем и удаляем все изображения из контента
        try:
            content = json.loads(article.content)
            for block in content:
                if block.get('type') == 'image' and block.get('path'):
                    file_path = os.path.join(settings.MEDIA_ROOT, block['path'])
                    if os.path.exists(file_path):
                        os.remove(file_path)
        except:
            pass  # Если не удалось загрузить контент, пропускаем

        # Удаляем все связанные ресурсы и их файлы
        for resource in article.resources.all():
            if resource.file:
                if os.path.exists(resource.file.path):
                    os.remove(resource.file.path)
            resource.delete()

        # Удаляем саму статью
        article.delete()
        
        messages.success(request, 'Статья и все связанные файлы успешно удалены')
        return redirect('teacher_articles_list')
    return render(request, 'main_app/delete_teacher_article.html', {'article': article})













